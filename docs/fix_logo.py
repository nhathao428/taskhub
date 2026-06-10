"""
Fix logo HUTECH:
- Slide 1: giữ ~2.5cm width (đã OK)
- Slide 17: resize từ 6.3×1.8cm (stretched) về 3.5×4.0cm (giữ aspect 415:480)
  + re-center theo title
- Bonus: cũng nhỏ logo slide 1 xuống 2cm nữa nếu cần
"""
from pptx import Presentation
from pptx.util import Emu, Cm

PPT = r'C:\Users\Admin\taskhub\docs\THUYET_TRINH_DO_AN.pptx'
LOGO_W_H_RATIO = 415 / 480  # 0.864 portrait

prs = Presentation(PPT)
slide_w = prs.slide_width
slide_h = prs.slide_height

# Slide 17 (idx 16): resize logo, center horizontally above title
s17 = prs.slides[16]
for shape in s17.shapes:
    if shape.shape_type == 13:  # PICTURE
        # New size: 3.5cm wide → height = 3.5/0.864 = 4.05cm
        new_w = Cm(3.5)
        new_h = Cm(4.05)
        shape.width = new_w
        shape.height = new_h
        # Center horizontally
        shape.left = (slide_w - new_w) // 2
        # Position: top area, ~1.5cm from top
        shape.top = Cm(1.2)
        print(f'Slide 17 logo: {Emu(shape.width).cm:.1f} × {Emu(shape.height).cm:.1f} cm at ({Emu(shape.left).cm:.1f}, {Emu(shape.top).cm:.1f})')

# Slide 1 (idx 0): keep but ensure proper aspect (was 2.4 × 2.0, slightly squashed)
s1 = prs.slides[0]
for shape in s1.shapes:
    if shape.shape_type == 13:
        new_w = Cm(2.0)
        new_h = Cm(2.32)  # 2.0 / 0.864
        shape.width = new_w
        shape.height = new_h
        print(f'Slide 1 logo: {Emu(shape.width).cm:.1f} × {Emu(shape.height).cm:.1f} cm at ({Emu(shape.left).cm:.1f}, {Emu(shape.top).cm:.1f})')

prs.save(PPT)
print('Saved')
