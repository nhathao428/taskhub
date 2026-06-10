import sys
import os
import shutil
from datetime import datetime
from PyQt5.QtWidgets import QApplication
from PyQt5.QtSvg import QSvgRenderer
from PyQt5.QtGui import QPainter, QImage
from PyQt5.QtCore import QSize, Qt
from pptx import Presentation
from pptx.util import Inches

def main():
    # 1. Initialize Qt Application for rendering
    app = QApplication(sys.argv)
    
    # 2. Setup paths
    base_dir = r"c:\Users\nhath\Downloads\taskhub\docs\ppt-project\thesis_defense_ppt169_20260530"
    svg_dir = os.path.join(base_dir, "svg_output")
    exports_dir = os.path.join(base_dir, "exports")
    temp_dir = os.path.join(base_dir, "temp_pngs")
    
    if not os.path.exists(svg_dir):
        print(f"Error: svg_output directory not found at {svg_dir}")
        return
        
    os.makedirs(exports_dir, exist_ok=True)
    os.makedirs(temp_dir, exist_ok=True)
    
    # 3. Find and sort SVG files
    svg_files = [f for f in os.listdir(svg_dir) if f.lower().endswith('.svg')]
    svg_files.sort()  # Will sort 01_cover.svg, 02_toc.svg, etc. in correct order
    
    if not svg_files:
        print("No SVG files found in svg_output!")
        return
        
    print(f"Found {len(svg_files)} SVG slides to convert.")
    
    # 4. Render SVG to PNG
    png_paths = []
    for svg_file in svg_files:
        svg_path = os.path.join(svg_dir, svg_file)
        png_name = svg_file[:-4] + ".png"
        png_path = os.path.join(temp_dir, png_name)
        
        print(f"Rendering {svg_file} to PNG...")
        renderer = QSvgRenderer(svg_path)
        if not renderer.isValid():
            print(f"  Warning: Invalid SVG file {svg_file}, skipping.")
            continue
            
        # Render at 1920x1080 (High resolution 16:9)
        image = QImage(1920, 1080, QImage.Format_ARGB32)
        image.fill(Qt.white)  # solid background
        
        painter = QPainter(image)
        renderer.render(painter)
        painter.end()
        
        if image.save(png_path):
            png_paths.append(png_path)
        else:
            print(f"  Warning: Failed to save PNG for {svg_file}")
            
    # 5. Create PowerPoint presentation
    print("\nCreating PPTX presentation...")
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank slide layout
    
    for png_path in png_paths:
        slide = prs.slides.add_slide(blank_layout)
        # Add the full-bleed image to the slide
        slide.shapes.add_picture(png_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        print(f"  Added slide: {os.path.basename(png_path)}")
        
    # 6. Save presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pptx_filename = f"thesis_defense_svg.pptx"
    pptx_timestamped_filename = f"thesis_defense_{timestamp}.pptx"
    
    pptx_path = os.path.join(exports_dir, pptx_filename)
    pptx_timestamped_path = os.path.join(exports_dir, pptx_timestamped_filename)
    
    prs.save(pptx_path)
    shutil.copy(pptx_path, pptx_timestamped_path)
    
    # Also copy to docs directory for convenience
    docs_pptx_path = r"c:\Users\nhath\Downloads\taskhub\docs\THUYET_TRINH_DO_AN_SVG.pptx"
    shutil.copy(pptx_path, docs_pptx_path)
    
    print("\n" + "="*50)
    print("SUCCESSFULLY CONVERTED PPT-PROJECT TO PPTX!")
    print(f"Saved to exports: {pptx_path}")
    print(f"Saved timestamped: {pptx_timestamped_path}")
    print(f"Saved copy to docs: {docs_pptx_path}")
    print("="*50)
    
    # 7. Clean up temporary PNGs
    try:
        shutil.rmtree(temp_dir)
        print("Temporary files cleaned up.")
    except Exception as e:
        print(f"Warning: Failed to clean up temp dir: {e}")

if __name__ == "__main__":
    main()
