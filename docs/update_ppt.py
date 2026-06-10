"""
Update PPT theo yêu cầu:
- Replace "TM" placeholder with HUTECH logo (slide 1 + slide 17)
- Add slide transitions (fade)
- Add text fade-in animation for body shapes
"""
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from lxml import etree

PPT = r'C:\Users\Admin\taskhub\docs\THUYET_TRINH_DO_AN.pptx'
LOGO = r'C:\Users\Admin\taskhub\docs\hutech_logo.png'

prs = Presentation(PPT)

# ============================================================
# 1. Replace "TM" placeholder with HUTECH logo
# ============================================================
replaced = 0
for slide in prs.slides:
    for shape in list(slide.shapes):
        if shape.has_text_frame and shape.text_frame.text.strip() == 'TM':
            # Capture position
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # Remove the text box
            sp = shape._element
            sp.getparent().remove(sp)
            # Add picture at same position (scale to fit width, keep aspect)
            slide.shapes.add_picture(LOGO, left, top, width=width, height=height)
            replaced += 1

print(f'Replaced TM with HUTECH logo: {replaced} slide(s)')

# ============================================================
# 2. Slide transitions — add <p:transition> to each slide XML
# ============================================================
NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def add_transition(slide_xml, kind='fade'):
    """Insert <p:transition spd="med"> after <p:cSld>."""
    # Remove existing transition if any
    for existing in slide_xml.findall('p:transition', NS):
        slide_xml.remove(existing)
    # Build transition element
    t = etree.SubElement(slide_xml, '{%s}transition' % NS['p'])
    t.set('spd', 'med')
    # Add the visual effect — fade is simple and clean
    if kind == 'fade':
        etree.SubElement(t, '{%s}fade' % NS['p'])
    elif kind == 'push':
        push = etree.SubElement(t, '{%s}push' % NS['p'])
        push.set('dir', 'l')
    return t


for slide in prs.slides:
    add_transition(slide._element, kind='fade')
print(f'Added fade transition to {len(prs.slides)} slides')

# ============================================================
# 3. Text fade-in animation — add timing/anim XML
# Use simple "fly in from left" entrance for body text shapes
# ============================================================
ANIM_TIMING_XML = """
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>__CHILDREN__</p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>
"""

ANIM_FADE_TEMPLATE = """
<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cTn id="__ID__" fill="hold">
    <p:stCondLst><p:cond delay="__DELAY__"/></p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="__ID2__" fill="hold">
          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="__ID3__" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="__ID4__" dur="1" fill="hold">
                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                      </p:cTn>
                      <p:tgtEl><p:spTgt spid="__SHAPEID__"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  <p:anim calcmode="lin" valueType="num">
                    <p:cBhvr additive="base">
                      <p:cTn id="__ID5__" dur="500" fill="hold"/>
                      <p:tgtEl><p:spTgt spid="__SHAPEID__"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:tavLst>
                      <p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
                      <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>
                    </p:tavLst>
                  </p:anim>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>
"""


def add_text_animations(slide):
    """Add fade-in entrance animation to text shapes (skipping title/header)."""
    sld_xml = slide._element
    # Remove existing timing if any
    for existing in sld_xml.findall('p:timing', NS):
        sld_xml.remove(existing)
    # Collect candidate shapes: text shapes with non-trivial text
    children_xml = []
    base_id = 3
    delay_ms = 0
    DELAY_STEP = 100  # 0.1s stagger
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        txt = shape.text_frame.text.strip()
        if not txt or len(txt) < 3: continue
        # Skip page numbers (single digit)
        if txt.isdigit() and len(txt) <= 2: continue
        sid = shape.shape_id
        ids = [base_id, base_id + 1, base_id + 2, base_id + 3, base_id + 4]
        base_id += 5
        block = ANIM_FADE_TEMPLATE
        block = block.replace('__SHAPEID__', str(sid))
        block = block.replace('__DELAY__', str(delay_ms))
        for j, k in enumerate(ids):
            block = block.replace(f'__ID{j+1 if j>0 else ""}__', str(k))
        # Plain __ID__ (no number) is first
        block = block.replace('__ID__', str(ids[0]))
        # Reset replace order: we used __ID__ for first which already done, need __ID2 __ID3 __ID4 __ID5
        children_xml.append(block)
        delay_ms += DELAY_STEP

    if not children_xml:
        return
    timing_str = ANIM_TIMING_XML.replace('__CHILDREN__', '\n'.join(children_xml))
    try:
        timing_el = etree.fromstring(timing_str)
        sld_xml.append(timing_el)
    except Exception as e:
        print(f'  anim XML parse error: {e}')


anim_count = 0
for slide in prs.slides:
    add_text_animations(slide)
    anim_count += 1
print(f'Added fade-in animations to {anim_count} slides')

prs.save(PPT)
print(f'\nSaved → {PPT}')
