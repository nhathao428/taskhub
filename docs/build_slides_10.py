# -*- coding: utf-8 -*-
"""
Slide thuyết trình ĐỒ ÁN CƠ SỞ — 10 slide trọng tâm, 16:9, có slide demo nhúng video.
Khớp yêu cầu: <=10 slide, trình bày ~8 phút.
Chạy: python build_slides_10.py  ->  THUYET_TRINH_DO_AN.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
UML = os.path.join(HERE, "uml", "png")
SHOTS = os.path.join(HERE, "screenshots")

# Palette
DEEPER = RGBColor(0x1E, 0x1B, 0x4B)
DEEP   = RGBColor(0x31, 0x2E, 0x81)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xEF, 0x44, 0x44)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xC7, 0xD2, 0xFE)
DARKTX = RGBColor(0x1F, 0x29, 0x37)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=DEEPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    # thanh accent trái
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = INDIGO; bar.line.fill.background()
    bar.shadow.inherit = False
    return s


def box(s, l, t, w, h):
    return s.shapes.add_textbox(l, t, w, h).text_frame


def txt(tf, text, size, color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT,
        space_after=6, font="Segoe UI"):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    return p


def bullet(tf, text, size=18, color=WHITE, lvl=0, bold=False, accent=None, space_after=10):
    p = tf.add_paragraph(); p.level = lvl; p.space_after = Pt(space_after)
    dot = p.add_run(); dot.text = "▸  "
    dot.font.size = Pt(size); dot.font.color.rgb = accent or CYAN; dot.font.bold = True; dot.font.name = "Segoe UI"
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = "Segoe UI"
    return p


def header(s, kicker, title, accent=INDIGO):
    tf = box(s, Inches(0.7), Inches(0.45), Inches(12), Inches(1.2))
    txt(tf, kicker.upper(), 14, accent if accent != INDIGO else CYAN, bold=True, space_after=2)
    txt(tf, title, 30, WHITE, bold=True, space_after=0)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.62), Inches(2.2), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = accent; ln.line.fill.background(); ln.shadow.inherit = False


def pageno(s, n):
    tf = box(s, Inches(12.2), Inches(6.95), Inches(1.0), Inches(0.4))
    txt(tf, f"{n}/10", 11, MUTED, align=PP_ALIGN.RIGHT, space_after=0)


def img(s, path, l, t, w=None, h=None):
    if os.path.exists(path):
        return s.shapes.add_picture(path, l, t, width=w, height=h)
    return None


def card(s, l, t, w, h, color, title, lines):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = DEEP; c.line.color.rgb = color; c.line.width = Pt(1.5)
    c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.14)
    txt(tf, title, 16, color, bold=True, space_after=6)
    for ln in lines:
        bullet(tf, ln, 13, WHITE, accent=color, space_after=5)
    return c


# ===================================================================== S1 BÌA
s = slide(DEEPER)
# vòng trang trí
o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.6), Inches(5.5), Inches(5.5))
o.fill.solid(); o.fill.fore_color.rgb = DEEP; o.line.fill.background(); o.shadow.inherit = False
img(s, os.path.join(HERE, "hutech_logo.png"), Inches(0.7), Inches(0.55), h=Inches(1.0))
tf = box(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.9))
txt(tf, "TRƯỜNG ĐẠI HỌC CÔNG NGHỆ TP.HCM  ·  KHOA CÔNG NGHỆ THÔNG TIN", 15, MUTED, bold=True, space_after=2)
txt(tf, "ĐỒ ÁN CƠ SỞ", 18, CYAN, bold=True, space_after=0)
tf = box(s, Inches(0.7), Inches(2.85), Inches(11.6), Inches(2.0))
txt(tf, "HỆ THỐNG QUẢN LÝ CÔNG VIỆC", 40, WHITE, bold=True, space_after=2)
txt(tf, "CHO DOANH NGHIỆP NHỎ ĐA NGÀNH TÍCH HỢP AI", 30, WHITE, bold=True, space_after=0)
tf = box(s, Inches(0.7), Inches(5.3), Inches(12), Inches(1.6))
txt(tf, "GVHD:  ThS. Dương Thành Phết", 18, MUTED, bold=True, space_after=4)
txt(tf, "SVTH:  Nguyễn Nhật Hào   ·   MSSV: 2380612688   ·   Lớp: 23DTHC1", 18, MUTED, bold=True, space_after=4)
txt(tf, "TP. Hồ Chí Minh, tháng 5 năm 2026", 14, MUTED, italic=True, space_after=0)

# ===================================================================== S2 ĐẶT VẤN ĐỀ
s = slide()
header(s, "1 · Đặt vấn đề", "Bài toán & Lý do chọn đề tài", RED)
tf = box(s, Inches(0.7), Inches(1.9), Inches(7.3), Inches(5))
bullet(tf, "DNNVV chiếm ~98% số doanh nghiệp Việt Nam, đóng góp 40–45% GDP", 19, accent=RED)
bullet(tf, "Phần lớn quản lý nhân sự & công việc thủ công: Excel, sổ sách, nhóm chat", 19, accent=RED)
bullet(tf, "Hệ quả: dữ liệu rời rạc, khó theo dõi tiến độ, dễ sai sót", 19, accent=RED)
bullet(tf, "Phân công nhân viên dựa trên cảm tính, thiếu cơ sở dữ liệu", 19, accent=RED)
card(s, Inches(8.3), Inches(2.2), Inches(4.3), Inches(3.2), CYAN, "Giải pháp đề xuất",
     ["Hệ thống quản lý công việc tập trung",
      "Phân quyền rõ ràng 3 vai trò",
      "AI gợi ý nhân viên phù hợp cho công việc",
      "Đa nền tảng: Web + Mobile"])
pageno(s, 2)

# ===================================================================== S3 MỤC TIÊU & PHẠM VI
s = slide()
header(s, "2 · Định hướng", "Mục tiêu & Phạm vi", CYAN)
card(s, Inches(0.7), Inches(2.0), Inches(5.9), Inches(4.3), CYAN, "Mục tiêu",
     ["Xác thực & phân quyền 3 vai trò (JWT)",
      "Quản lý nhân viên · dự án · công việc",
      "Chấm công vào/ra theo ngày",
      "AI gợi ý nhân viên (top 5 + lý do)",
      "Dashboard thống kê tổng quan"])
card(s, Inches(6.9), Inches(2.0), Inches(5.7), Inches(4.3), AMBER, "Phạm vi đồ án cơ sở",
     ["Tập trung phần LÕI + module AI gợi ý",
      "Tính năng nâng cao để dành Đồ án chuyên ngành:",
      "   chấm công GPS · cache Redis",
      "   ứng dụng mobile · triển khai cloud"])
pageno(s, 3)

# ===================================================================== S4 CÔNG NGHỆ & KIẾN TRÚC
s = slide()
header(s, "3 · Cơ sở lý thuyết", "Công nghệ & Kiến trúc hệ thống", GREEN)
img(s, os.path.join(UML, "architecture.png"), Inches(0.7), Inches(2.0), w=Inches(7.2))
card(s, Inches(8.3), Inches(2.0), Inches(4.3), Inches(4.3), GREEN, "Kiến trúc 3 lớp",
     ["Client: React + Vite + Tailwind",
      "Application: Spring Boot + JWT",
      "Data: PostgreSQL 16",
      "AI: Google Gemini 2.5 Flash",
      "REST API · Spring Security · JPA"])
pageno(s, 4)

# ===================================================================== S5 PHÂN TÍCH & THIẾT KẾ
s = slide()
header(s, "4 · Phân tích – Thiết kế", "Use Case & Cơ sở dữ liệu", INDIGO)
img(s, os.path.join(UML, "use-case-tong-the.png"), Inches(0.7), Inches(2.0), h=Inches(4.5))
card(s, Inches(6.2), Inches(2.0), Inches(6.4), Inches(4.5), CYAN, "Thiết kế",
     ["3 tác nhân: Admin · Manager · Employee",
      "14 use case chính",
      "CSDL 7 bảng (ERD)",
      "users · employees · projects",
      "tasks · attendance · suggestions"])
pageno(s, 5)

# ===================================================================== S6 CHỨC NĂNG THEO VAI TRÒ
s = slide()
header(s, "5 · Chức năng", "Phân quyền theo 3 vai trò", AMBER)
card(s, Inches(0.7), Inches(2.0), Inches(3.85), Inches(4.3), RED, "ADMIN",
     ["Toàn quyền hệ thống", "Quản lý người dùng", "Phân quyền tài khoản", "Cấu hình hệ thống"])
card(s, Inches(4.75), Inches(2.0), Inches(3.85), Inches(4.3), INDIGO, "MANAGER",
     ["Quản lý nhân viên, dự án", "Tạo & gán công việc", "Theo dõi chấm công", "AI gợi ý nhân viên", "Dashboard thống kê"])
card(s, Inches(8.8), Inches(2.0), Inches(3.8), Inches(4.3), GREEN, "EMPLOYEE",
     ["Xem công việc được giao", "Cập nhật trạng thái", "Tự check-in / check-out", "Xem lịch sử chấm công"])
pageno(s, 6)

# ===================================================================== S7 ĐIỂM NHẤN: AI
s = slide()
header(s, "6 · Điểm nhấn", "Module AI gợi ý nhân viên (Google Gemini)", RGBColor(0x93, 0x33, 0xEA))
tf = box(s, Inches(0.7), Inches(1.95), Inches(7.4), Inches(5))
bullet(tf, "Đầu vào: công việc cần phân công + dữ liệu hiệu suất nhân viên", 18, accent=CYAN)
bullet(tf, "Dữ liệu: tiến độ task · tỷ lệ đúng hạn · chấm công · kỹ năng", 18, accent=CYAN)
bullet(tf, "Xử lý: dựng prompt tiếng Việt → gửi Gemini xếp hạng", 18, accent=CYAN)
bullet(tf, "Backend KHÔNG tính điểm cứng — AI tự đánh giá & giải thích", 18, accent=AMBER, bold=True)
bullet(tf, "Đầu ra: TOP 5 nhân viên phù hợp + lý do bằng tiếng Việt", 18, accent=GREEN, bold=True)
img(s, os.path.join(SHOTS, "09_ai_result.png"), Inches(8.4), Inches(2.15), w=Inches(4.2))
pageno(s, 7)

# ===================================================================== S8 DEMO (LIVE)
s = slide(DEEPER)
header(s, "7 · Sản phẩm", "Trình diễn trực tiếp hệ thống (Demo)", GREEN)
tf = box(s, Inches(0.7), Inches(2.8), Inches(12), Inches(3.0))
txt(tf, "DEMO HỆ THỐNG TRỰC TIẾP", 36, CYAN, bold=True, align=PP_ALIGN.CENTER, space_after=16)
txt(tf, "Trình diễn trực tiếp thao tác trên Web (Quản lý) và giao diện Mobile (Nhân viên)", 18, WHITE, align=PP_ALIGN.CENTER, space_after=10)
txt(tf, "· Quy trình Đăng nhập, Thêm nhân viên, Dự án, Công việc và Chấm công\n· Quy trình chạy AI gợi ý nhân sự tối ưu hiệu suất công việc", 16, MUTED, align=PP_ALIGN.CENTER, space_after=0)
pageno(s, 8)

# ===================================================================== S9 KIỂM THỬ
s = slide()
header(s, "8 · Kiểm thử", "Kiểm thử & Đánh giá kết quả", CYAN)
tf = box(s, Inches(0.7), Inches(1.95), Inches(7.0), Inches(5))
bullet(tf, "Phương pháp: kiểm thử hộp đen ở mức chức năng", 19)
bullet(tf, "Bao phủ 5 module: Auth · Nhân viên · Dự án–Công việc · Chấm công · AI", 19)
bullet(tf, "32 ca kiểm thử → đạt 100%", 19, accent=GREEN, bold=True)
bullet(tf, "Hệ thống ổn định, phân quyền chặt, AI gợi ý hợp lý", 19, accent=GREEN)
card(s, Inches(8.1), Inches(2.1), Inches(4.5), Inches(3.4), GREEN, "Kết quả kiểm thử",
     ["Xác thực – phân quyền: 8/8",
      "Quản lý nhân viên: 6/6",
      "Dự án – công việc: 8/8",
      "Chấm công: 5/5",
      "AI gợi ý: 5/5   →   Tổng 32/32"])
pageno(s, 9)

# ===================================================================== S10 KẾT LUẬN
s = slide()
header(s, "9 · Kết luận", "Kết luận & Hướng phát triển", INDIGO)
card(s, Inches(0.7), Inches(2.0), Inches(5.9), Inches(4.3), GREEN, "Kết quả đạt được",
     ["Ứng dụng web 3 lớp hoàn chỉnh",
      "Xác thực & phân quyền JWT 3 vai trò",
      "CRUD nhân viên · dự án · công việc · chấm công",
      "Module AI gợi ý nhân viên (Gemini)"])
card(s, Inches(6.9), Inches(2.0), Inches(5.7), Inches(4.3), AMBER, "Hướng phát triển",
     ["Chấm công GPS (geofence)",
      "Cache Redis · giới hạn truy vấn",
      "Ứng dụng mobile đa nền tảng",
      "Multi-tenant + gói thuê bao (thương mại hóa)"])
tf = box(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.7))
txt(tf, "Xin chân thành cảm ơn quý thầy cô đã lắng nghe — Rất mong nhận được câu hỏi & góp ý!",
    16, CYAN, bold=True, align=PP_ALIGN.CENTER, space_after=0)
pageno(s, 10)

OUT = os.path.join(HERE, "THUYET_TRINH_DO_AN.pptx")
prs.save(OUT)
print(f"[OK] {len(prs.slides._sldIdLst)} slide -> {OUT}")
