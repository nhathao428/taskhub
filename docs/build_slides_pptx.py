# -*- coding: utf-8 -*-
"""
Sinh slide thuyết trình đồ án (.pptx) — 17 slide, 16:9.
Chạy: python build_slides_pptx.py  ->  THUYET_TRINH_DO_AN.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))

# ---- Bảng màu ----
INDIGO  = RGBColor(0x4F, 0x46, 0xE5)
DEEP    = RGBColor(0x31, 0x2E, 0x81)
DEEPER  = RGBColor(0x1E, 0x1B, 0x4B)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
PINK    = RGBColor(0xEC, 0x48, 0x99)
CYAN    = RGBColor(0x06, 0xB6, 0xD4)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x1E, 0x29, 0x3B)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
FONT    = "Arial"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, content, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
         line_spacing=1.1):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = content if isinstance(content, list) else [content]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def bullets(s, x, y, w, h, items, size=18, color=INK, gap=10,
            glyph_color=INDIGO):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        g = p.add_run()
        g.text = ("▸  " if lvl == 0 else "–  ")
        g.font.name = FONT
        g.font.size = Pt(size)
        g.font.bold = True
        g.font.color.rgb = glyph_color if lvl == 0 else MUTED
        r = p.add_run()
        r.text = txt
        r.font.name = FONT
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = color
        if lvl == 1:
            p.level = 1
    return tb


def pic_fit(s, path, cx, cy, max_w, max_h):
    p = os.path.join(HERE, path)
    pic = s.shapes.add_picture(p, 0, 0, width=Inches(max_w))
    if pic.height > Inches(max_h):
        ratio = Inches(max_h) / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = Inches(max_h)
    pic.left = int(Inches(cx) - pic.width / 2)
    pic.top = int(Inches(cy) - pic.height / 2)
    return pic


def transition(s):
    """Thêm hiệu ứng chuyển slide (fade)."""
    el = s._element
    tr = el.find(qn('p:transition'))
    if tr is None:
        tr = etree.SubElement(el, qn('p:transition'))
    tr.set('spd', 'med')
    for c in list(tr):
        tr.remove(c)
    etree.SubElement(tr, qn('p:fade'))


def header(s, kicker, title, idx):
    """Thanh tiêu đề cho slide nội dung."""
    bg(s, WHITE)
    rect(s, 0, 0, 13.333, 1.32, INDIGO)
    rect(s, 0, 1.32, 13.333, 0.06, PINK)
    rect(s, 0.55, 0.30, 0.16, 0.72, AMBER)
    text(s, 0.95, 0.22, 11.5, 0.4, kicker.upper(), size=12, color=RGBColor(0xC7,0xD2,0xFE), bold=True)
    text(s, 0.95, 0.50, 11.5, 0.7, title, size=27, color=WHITE, bold=True)
    # footer
    text(s, 0.55, 7.02, 9.0, 0.4,
         "Hệ thống Quản lý Công việc cho Doanh nghiệp Nhỏ Tích hợp AI",
         size=9, color=MUTED)
    text(s, 12.0, 7.02, 0.8, 0.4, str(idx), size=11, color=INDIGO, bold=True,
         align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, label, value, accent):
    """Ô thông tin nhỏ."""
    rect(s, x, y, w, 1.5, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, 0.12, 1.5, accent)
    text(s, x + 0.28, y + 0.16, w - 0.4, 0.4, label.upper(), size=11,
         color=MUTED, bold=True)
    text(s, x + 0.28, y + 0.50, w - 0.4, 0.9, value, size=15, color=INK,
         bold=True)


# ============================================================
# SLIDE 1 — Bìa
# ============================================================
s = slide()
bg(s, DEEPER)
rect(s, 0, 0, 5.0, 7.5, DEEP)
rect(s, 0, 0, 0.22, 7.5, PINK)
# khối logo
rect(s, 0.95, 0.85, 0.95, 0.95, INDIGO, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 0.95, 0.95, 0.95, 0.8, "TM", size=30, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, 0.6, 2.1, 4.0, 2.5,
     ["BỘ GIÁO DỤC VÀ ĐÀO TẠO", "TRƯỜNG ĐẠI HỌC CÔNG NGHỆ TP.HCM",
      "KHOA CÔNG NGHỆ THÔNG TIN"],
     size=12, color=RGBColor(0xC7,0xD2,0xFE), bold=True, line_spacing=1.5)
text(s, 0.6, 5.7, 4.0, 0.5, "ĐỒ ÁN CƠ SỞ", size=15, color=AMBER, bold=True)
text(s, 0.6, 6.1, 4.0, 0.5, "Tháng 5 năm 2026", size=12,
     color=RGBColor(0xA5,0xB4,0xFC))

text(s, 5.7, 1.55, 7.2, 0.5, "ĐỀ TÀI", size=14, color=PINK, bold=True)
text(s, 5.65, 1.95, 7.4, 2.6,
     "HỆ THỐNG QUẢN LÝ CÔNG VIỆC CHO DOANH NGHIỆP NHỎ ĐA NGÀNH TÍCH HỢP AI",
     size=33, color=WHITE, bold=True, line_spacing=1.12)
rect(s, 5.7, 4.35, 1.7, 0.07, PINK)
text(s, 5.7, 4.55, 7.2, 0.5,
     "Ứng dụng Web · Mobile · tích hợp Google Gemini AI",
     size=15, color=RGBColor(0xC7,0xD2,0xFE), italic=True)

info = [
    ("Giảng viên hướng dẫn", "ThS. Dương Thành Phết"),
    ("Sinh viên thực hiện", "Nguyễn Nhật Hảo"),
    ("Mã số sinh viên", "2380612688"),
    ("Lớp", "23DTHC1"),
]
yy = 5.25
for lb, vl in info:
    text(s, 5.7, yy, 3.0, 0.4, lb, size=12, color=RGBColor(0xA5,0xB4,0xFC), bold=True)
    text(s, 8.7, yy, 4.2, 0.4, vl, size=13, color=WHITE, bold=True)
    yy += 0.46
transition(s)

# ============================================================
# SLIDE 2 — Nội dung trình bày
# ============================================================
s = slide()
header(s, "Mở đầu", "Nội dung trình bày", 2)
agenda = [
    ("1", "Đặt vấn đề & mục tiêu", INDIGO),
    ("2", "Công nghệ & kiến trúc hệ thống", PURPLE),
    ("3", "Các chức năng chính", PINK),
    ("4", "Điểm nhấn: AI gợi ý & giao diện song ngữ", CYAN),
    ("5", "Kiểm thử & triển khai trên AWS", GREEN),
    ("6", "Kết quả, hạn chế & hướng phát triển", AMBER),
]
yy = 1.85
for num, label, c in agenda:
    rect(s, 0.95, yy, 0.62, 0.62, c, shape=MSO_SHAPE.OVAL)
    text(s, 0.95, yy + 0.04, 0.62, 0.55, num, size=20, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, 1.85, yy + 0.07, 10.5, 0.55, label, size=21, color=INK, bold=True)
    yy += 0.83
transition(s)

# ============================================================
# SLIDE 3 — Đặt vấn đề
# ============================================================
s = slide()
header(s, "Phần 1 · Tổng quan", "Đặt vấn đề", 3)
bullets(s, 0.95, 1.75, 11.4, 5.0, [
    "Doanh nghiệp nhỏ đa ngành quản lý công việc thủ công qua Excel, tin nhắn, giấy tờ — dữ liệu phân tán, khó theo dõi.",
    "Không có nguồn thông tin thống nhất (single source of truth) về tiến độ công việc.",
    "Phân công nhân sự theo cảm tính, chưa dựa trên dữ liệu năng lực và lịch sử thực tế.",
    "Phần mềm thương mại đắt, nặng, nhiều tính năng thừa so với nhu cầu doanh nghiệp nhỏ.",
    "Chưa khai thác trí tuệ nhân tạo (AI) để hỗ trợ ra quyết định phân công.",
], size=19, gap=16)
transition(s)

# ============================================================
# SLIDE 4 — Mục tiêu đề tài
# ============================================================
s = slide()
header(s, "Phần 1 · Tổng quan", "Mục tiêu đề tài", 4)
bullets(s, 0.95, 1.75, 11.4, 5.0, [
    "Xây dựng hệ thống quản lý công việc gọn nhẹ, chi phí thấp, phù hợp doanh nghiệp nhỏ.",
    "Kiến trúc đa nền tảng: Web (React) + Mobile (Flutter) + Backend (Spring Boot).",
    "Tích hợp AI (Google Gemini) gợi ý nhân viên phù hợp nhất cho từng công việc.",
    "Bảo mật bằng JWT, phân quyền 3 vai trò: ADMIN / MANAGER / EMPLOYEE.",
    "Đóng gói bằng Docker, triển khai được thực tế trên hạ tầng đám mây.",
], size=19, gap=16)
transition(s)

# ============================================================
# SLIDE 5 — Công nghệ sử dụng
# ============================================================
s = slide()
header(s, "Phần 2 · Công nghệ", "Công nghệ sử dụng", 5)
tech = [
    ("Backend", "Java 17 · Spring Boot 3.5\nSpring Security + JWT · JPA", INDIGO),
    ("Frontend", "React 18 · Vite\nTailwind CSS · Chart.js", PURPLE),
    ("Mobile", "Flutter 3.x · Dart\nchạy đa nền tảng", PINK),
    ("CSDL & Cache", "PostgreSQL 16\nRedis 7", CYAN),
    ("AI", "Google Gemini\ngemini-2.5-flash", GREEN),
    ("Hạ tầng", "Docker · Caddy\nAWS EC2", AMBER),
]
x0, y0, cw, ch, gx, gy = 0.95, 1.85, 3.74, 2.35, 0.18, 0.30
for i, (name, desc, c) in enumerate(tech):
    cx = x0 + (i % 3) * (cw + gx)
    cy = y0 + (i // 3) * (ch + gy)
    rect(s, cx, cy, cw, ch, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, cy, cw, 0.12, c)
    text(s, cx + 0.3, cy + 0.28, cw - 0.6, 0.5, name, size=17, color=c, bold=True)
    text(s, cx + 0.3, cy + 0.82, cw - 0.6, 1.4, desc.split("\n"), size=13,
         color=INK, line_spacing=1.3)
transition(s)

# ============================================================
# SLIDE 6 — Kiến trúc hệ thống
# ============================================================
s = slide()
header(s, "Phần 2 · Công nghệ", "Kiến trúc tổng thể hệ thống", 6)
pic_fit(s, "uml/png/architecture.png", 4.05, 4.35, 6.7, 4.85)
bullets(s, 8.1, 2.0, 4.7, 4.7, [
    "Mô hình 3 tầng: Client – Ứng dụng – Dữ liệu.",
    "Client: Web React + Mobile Flutter.",
    "Tầng ứng dụng: Spring Boot REST API, bảo mật JWT, cache.",
    "Tầng dữ liệu: PostgreSQL + Redis.",
    "AI qua Google Gemini API.",
], size=15, gap=12)
transition(s)

# ============================================================
# SLIDE 7 — Sơ đồ Use Case
# ============================================================
s = slide()
header(s, "Phần 2 · Phân tích thiết kế", "Sơ đồ Use Case tổng thể", 7)
pic_fit(s, "uml/png/use-case-tong-the.png", 6.66, 4.4, 12.2, 5.1)
transition(s)

# ============================================================
# SLIDE 8 — Các chức năng chính
# ============================================================
s = slide()
header(s, "Phần 3 · Chức năng", "Các chức năng chính", 8)
funcs = [
    ("Xác thực & phân quyền", "3 vai trò, JWT", INDIGO),
    ("Quản lý nhân viên", "CRUD kèm kỹ năng", PURPLE),
    ("Quản lý dự án", "CRUD, gán nhân sự", PINK),
    ("Quản lý công việc", "Giao việc, tiến độ", CYAN),
    ("Chấm công GPS", "Geofence đa văn phòng", GREEN),
    ("AI gợi ý nhân viên", "Top 5 phù hợp", AMBER),
    ("Dashboard thống kê", "Biểu đồ trực quan", INDIGO),
    ("Giao diện song ngữ", "Việt / Anh", PURPLE),
]
x0, y0, cw, ch, gx, gy = 0.95, 1.85, 2.78, 2.3, 0.18, 0.28
for i, (name, desc, c) in enumerate(funcs):
    cx = x0 + (i % 4) * (cw + gx)
    cy = y0 + (i // 4) * (ch + gy)
    rect(s, cx, cy, cw, ch, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         line=RGBColor(0xE2,0xE8,0xF0))
    rect(s, cx + 0.28, cy + 0.3, 0.5, 0.5, c, shape=MSO_SHAPE.OVAL)
    text(s, cx + 0.26, cy + 0.95, cw - 0.5, 0.85, name, size=14, color=INK, bold=True)
    text(s, cx + 0.26, cy + 1.62, cw - 0.5, 0.5, desc, size=11, color=MUTED)
transition(s)

# ============================================================
# SLIDE 9 — AI gợi ý nhân viên
# ============================================================
s = slide()
header(s, "Phần 4 · Điểm nhấn", "AI gợi ý nhân viên phù hợp", 9)
steps = [
    ("1", "Nhập yêu cầu", "Manager nhập tiêu đề, mô tả công việc, kỹ năng yêu cầu", INDIGO),
    ("2", "Gom dữ liệu", "Backend tổng hợp lịch sử task, đúng hạn, chấm công, kỹ năng", PURPLE),
    ("3", "Gọi Gemini", "Xây prompt → Google Gemini gemini-2.5-flash xếp hạng", PINK),
    ("4", "Trả kết quả", "TOP 5 nhân viên kèm lý do; cache Redis 5 phút", GREEN),
]
yy = 1.9
for num, t1, t2, c in steps:
    rect(s, 0.95, yy, 0.7, 0.7, c, shape=MSO_SHAPE.OVAL)
    text(s, 0.95, yy + 0.07, 0.7, 0.6, num, size=22, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    text(s, 1.95, yy + 0.02, 3.2, 0.6, t1, size=18, color=c, bold=True)
    text(s, 1.95, yy + 0.42, 10.6, 0.6, t2, size=14, color=INK)
    yy += 1.0
rect(s, 0.95, 6.05, 11.4, 0.92, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 1.25, 6.22, 11.0, 0.6,
     "Cache Redis: request lặp lại trên cùng công việc giảm từ ~1450 ms xuống dưới 10 ms.",
     size=14, color=INK, bold=True)
transition(s)

# ============================================================
# SLIDE 10 — Giao diện song ngữ
# ============================================================
s = slide()
header(s, "Phần 4 · Điểm nhấn", "Giao diện song ngữ Việt / Anh", 10)
bullets(s, 0.95, 1.85, 6.6, 4.8, [
    "Toàn bộ giao diện hỗ trợ Tiếng Việt và English, chuyển đổi tức thì.",
    "Cơ chế i18n tự xây dựng: LanguageContext + từ điển dịch + hàm t().",
    "Nút cờ VI / EN, lưu lựa chọn vào localStorage — nhớ ở lần sau.",
    "Dịch hoàn toàn phía client, KHÔNG phụ thuộc Google Translate.",
    "Mọi trang và component đều lấy chuỗi hiển thị qua hàm t().",
], size=16, gap=14)
pic_fit(s, "screenshots/01_login.png", 10.15, 4.35, 5.3, 4.7)
transition(s)

# ============================================================
# SLIDE 11 — Chấm công GPS
# ============================================================
s = slide()
header(s, "Phần 3 · Chức năng", "Chấm công GPS & Geofence", 11)
bullets(s, 0.95, 1.85, 6.5, 4.8, [
    "Cấu hình nhiều văn phòng: toạ độ + bán kính cho phép.",
    "Chấm công kèm vị trí GPS, tính khoảng cách Haversine.",
    "Ngoài vùng hoặc nghi GPS giả → trạng thái PENDING_REVIEW chờ quản lý duyệt.",
    "Bản đồ Leaflet / OpenStreetMap hiển thị trực quan.",
], size=16, gap=14)
pic_fit(s, "screenshots/15_my_attendance_map.png", 10.1, 4.35, 5.4, 4.7)
transition(s)

# ============================================================
# SLIDE 12 — Demo giao diện
# ============================================================
s = slide()
header(s, "Phần 3 · Sản phẩm", "Video demo hệ thống", 12)
s.shapes.add_movie(
    os.path.join(HERE, "demo_video.mp4"),
    Inches(2.07), Inches(1.66), Inches(9.2), Inches(5.18),
    poster_frame_image=os.path.join(HERE, "demo_poster.png"),
    mime_type="video/mp4",
)
text(s, 0.95, 6.95, 11.4, 0.4,
     "Video tham quan các màn hình chính — bấm để phát khi thuyết trình",
     size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
transition(s)

# ============================================================
# SLIDE 13 — Kiểm thử
# ============================================================
s = slide()
header(s, "Phần 5 · Kiểm thử", "Kiểm thử & đánh giá", 13)
chip(s, 0.95, 1.9, 3.62, "Unit test backend", "34 test · 100% pass", GREEN)
chip(s, 4.85, 1.9, 3.62, "Build production", "Backend + Frontend sạch", INDIGO)
chip(s, 8.75, 1.9, 3.62, "Nền tảng", "Web + Mobile thật", PURPLE)
bullets(s, 0.95, 3.85, 11.4, 3.0, [
    "Kiểm thử hộp đen theo từng module: Auth, Employee, Task, Attendance, AI.",
    "Unit test bằng JUnit 5 + Mockito cho tầng service và bảo mật.",
    "Kiểm thử thực nghiệm trên trình duyệt thật và thiết bị di động.",
    "Đo hiệu năng API AI: xác nhận hiệu quả của cache Redis.",
], size=17, gap=14)
transition(s)

# ============================================================
# SLIDE 14 — Triển khai AWS
# ============================================================
s = slide()
header(s, "Phần 5 · Triển khai", "Triển khai thực tế trên AWS", 14)
chip(s, 0.95, 1.9, 3.62, "Máy chủ", "AWS EC2 t3.small\nUbuntu 24.04 LTS", INDIGO)
chip(s, 4.85, 1.9, 3.62, "Đóng gói", "Docker Compose\n5 container", PURPLE)
chip(s, 8.75, 1.9, 3.62, "Truy cập", "Công khai qua\nElastic IP", GREEN)
bullets(s, 0.95, 3.95, 11.4, 2.9, [
    "Toàn bộ hệ thống chạy qua Docker Compose: PostgreSQL + Redis + backend + frontend + Caddy.",
    "Caddy reverse proxy: định tuyến /api tới backend, còn lại tới frontend.",
    "Tự khởi động lại sau khi server reboot (restart: unless-stopped).",
    "Sản phẩm vận hành thực tế trên Internet, không chỉ chạy cục bộ.",
], size=16, gap=13)
transition(s)

# ============================================================
# SLIDE 15 — Kết quả đạt được
# ============================================================
s = slide()
header(s, "Phần 6 · Kết luận", "Kết quả đạt được", 15)
bullets(s, 0.95, 1.8, 11.4, 5.0, [
    "Hoàn thành 14/14 yêu cầu chức năng đã đặc tả (đạt 100%).",
    "Gần 30 REST endpoint, phân quyền hai lớp theo 3 vai trò.",
    "Đầy đủ Web + Mobile + AI gợi ý + giao diện song ngữ + chấm công GPS.",
    "Đã triển khai chạy thật trên hạ tầng đám mây AWS EC2.",
    "Tài liệu đầy đủ: báo cáo đồ án, sơ đồ UML, hướng dẫn cài đặt & triển khai.",
], size=19, gap=16)
transition(s)

# ============================================================
# SLIDE 16 — Hạn chế & hướng phát triển
# ============================================================
s = slide()
header(s, "Phần 6 · Kết luận", "Hạn chế & hướng phát triển", 16)
rect(s, 0.95, 1.8, 5.6, 0.55, RGBColor(0xFE,0xE2,0xE2),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 1.2, 1.86, 5.2, 0.45, "HẠN CHẾ", size=14, color=RGBColor(0xB9,0x1C,0x1C), bold=True)
bullets(s, 0.95, 2.5, 5.7, 4.3, [
    "Phụ thuộc Gemini API (giới hạn free tier).",
    "Chưa có thông báo real-time.",
    "Chưa có refresh token cho JWT.",
    "Độ phủ unit test backend còn thấp (~15%).",
], size=14, gap=12, glyph_color=RGBColor(0xB9,0x1C,0x1C))
rect(s, 6.95, 1.8, 5.4, 0.55, RGBColor(0xDC,0xFC,0xE7),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 7.2, 1.86, 5.0, 0.45, "HƯỚNG PHÁT TRIỂN", size=14, color=RGBColor(0x15,0x80,0x3D), bold=True)
bullets(s, 6.95, 2.5, 5.5, 4.3, [
    "WebSocket cho thông báo real-time.",
    "Refresh token + thu hồi token.",
    "RAG + embeddings nâng cao AI gợi ý.",
    "CI/CD tự động; mở rộng sang Kubernetes.",
    "Phát triển phiên bản SaaS đa doanh nghiệp.",
], size=14, gap=12, glyph_color=GREEN)
transition(s)

# ============================================================
# SLIDE 17 — Cảm ơn
# ============================================================
s = slide()
bg(s, DEEPER)
rect(s, 0, 0, 13.333, 0.22, PINK)
rect(s, 0, 7.28, 13.333, 0.22, PINK)
rect(s, 5.42, 1.75, 2.5, 0.95, INDIGO, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 5.42, 1.9, 2.5, 0.7, "TM", size=34, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER)
text(s, 1.5, 3.0, 10.3, 1.2, "CẢM ƠN THẦY CÔ ĐÃ LẮNG NGHE", size=40,
     color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, 1.5, 4.3, 10.3, 0.6, "Rất mong nhận được câu hỏi và góp ý từ hội đồng",
     size=17, color=RGBColor(0xC7,0xD2,0xFE), align=PP_ALIGN.CENTER, italic=True)
text(s, 1.5, 5.5, 10.3, 0.5, "Nguyễn Nhật Hảo  ·  23DTHC1  ·  GVHD: ThS. Dương Thành Phết",
     size=14, color=RGBColor(0xA5,0xB4,0xFC), align=PP_ALIGN.CENTER)
transition(s)

out = os.path.join(HERE, "THUYET_TRINH_DO_AN.pptx")
prs.save(out)
print("[OK] Saved:", out, "—", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
