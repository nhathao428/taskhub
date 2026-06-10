# -*- coding: utf-8 -*-
"""Tạo lại THUYET_TRINH_DO_AN.pptx từ 6 trang brochure JPEG + hiệu ứng.

- Backup pptx hiện tại trước khi ghi đè.
- 6 slide 16:9, mỗi slide 1 ảnh full-bleed.
- Transition fade giữa các slide + animation fade-in cho ảnh.
"""
import os
import shutil
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches
from lxml import etree

DOCS = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BRO = os.path.join(DOCS, "brochure")
PPT = os.path.join(DOCS, "THUYET_TRINH_DO_AN.pptx")

PAGES = [
    "TaskHub_gioithieu_01_bia.jpg",
    "TaskHub_gioithieu_02_baitoan.jpg",
    "TaskHub_gioithieu_03_tinhnang.jpg",
    "TaskHub_gioithieu_04_AI.jpg",
    "TaskHub_gioithieu_05_kientruc.jpg",
    "TaskHub_gioithieu_06_doituong.jpg",
]

NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

# 1) Backup bản hiện tại
if os.path.exists(PPT):
    stamp = datetime.now().strftime("%Y%m%d-%H%M")
    bak = os.path.join(DOCS, f"THUYET_TRINH_DO_AN.backup-before-brochure-{stamp}.pptx")
    shutil.copy(PPT, bak)
    print("Backup ->", os.path.basename(bak))

# 2) Tạo presentation mới 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # layout trống

pic_ids = []
for fname in PAGES:
    slide = prs.slides.add_slide(blank)
    pic = slide.shapes.add_picture(
        os.path.join(BRO, fname), 0, 0,
        width=prs.slide_width, height=prs.slide_height,
    )
    pic_ids.append((slide, pic.shape_id))

# 3) Transition fade cho từng slide
def add_transition(slide_el):
    for ex in slide_el.findall("p:transition", NS):
        slide_el.remove(ex)
    t = etree.SubElement(slide_el, "{%s}transition" % NS["p"])
    t.set("spd", "med")
    etree.SubElement(t, "{%s}fade" % NS["p"])

# 4) Animation fade-in (entrance) cho ảnh mỗi slide
TIMING = """
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
    <p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
      <p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
        <p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
          <p:par><p:cTn id="5" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
              <p:set><p:cBhvr><p:cTn id="6" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                <p:tgtEl><p:spTgt spid="__SID__"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>
                <p:to><p:strVal val="visible"/></p:to></p:set>
              <p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base">
                <p:cTn id="7" dur="600" fill="hold"/><p:tgtEl><p:spTgt spid="__SID__"/></p:tgtEl>
                <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst></p:cBhvr>
                <p:tavLst><p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
                <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav></p:tavLst></p:anim>
            </p:childTnLst></p:cTn></p:par>
        </p:childTnLst></p:cTn></p:par>
      </p:childTnLst></p:cTn></p:par>
    </p:childTnLst></p:cTn>
    <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
    <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
    </p:seq>
  </p:childTnLst></p:cTn></p:par></p:tnLst>
</p:timing>
"""

for slide, sid in pic_ids:
    sld_el = slide._element
    add_transition(sld_el)
    for ex in sld_el.findall("p:timing", NS):
        sld_el.remove(ex)
    sld_el.append(etree.fromstring(TIMING.replace("__SID__", str(sid))))

prs.save(PPT)
print(f"Saved -> {PPT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
