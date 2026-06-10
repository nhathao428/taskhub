# -*- coding: utf-8 -*-
"""Tạo THUYET_TRINH_DO_AN.pptx từ 7 ảnh brochure full-bleed (pixel-perfect mọi app)
+ nhúng video demo native phủ đúng khung poster ở slide 6. Transition fade + ảnh fade-in."""
import os, shutil
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

HOME = r"C:\Users\Admin"
DOCS = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BRO = os.path.join(DOCS, "brochure")
PPT = os.path.join(DOCS, "THUYET_TRINH_DO_AN.pptx")
NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

NAMES = [
    "TaskHub_gioithieu_01_bia.jpg",
    "TaskHub_gioithieu_02_baitoan.jpg",
    "TaskHub_gioithieu_03_tinhnang.jpg",
    "TaskHub_gioithieu_04_AI.jpg",
    "TaskHub_gioithieu_05_kientruc.jpg",
    "TaskHub_gioithieu_06_demo.jpg",
    "TaskHub_gioithieu_07_doituong.jpg",
]
# 1) copy ảnh render (n1..n7) vào brochure với tên chuẩn
for i, nm in enumerate(NAMES, 1):
    src = os.path.join(HOME, f"n{i}.jpeg")
    if os.path.exists(src):
        shutil.copy(src, os.path.join(BRO, nm))
print("Đã cập nhật 7 ảnh brochure")

# 2) backup pptx
if os.path.exists(PPT):
    stamp = datetime.now().strftime("%Y%m%d-%H%M")
    bak = os.path.join(DOCS, f"THUYET_TRINH_DO_AN.backup-{stamp}.pptx")
    shutil.copy(PPT, bak); print("Backup ->", os.path.basename(bak))

# 3) deck 16:9, mỗi slide 1 ảnh full-bleed
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

PX = 13.333 / 1920.0  # inch / pixel
DEMO_IDX = 5  # slide thứ 6 (0-based)
# vùng poster trong ảnh demo (px): left=735 top=358 w=1100 h=642
vid = dict(x=735 * PX, y=358 * PX, w=1100 * PX, h=642 * PX)

pic_shapes = []
for idx, nm in enumerate(NAMES):
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(os.path.join(BRO, nm), 0, 0, width=prs.slide_width, height=prs.slide_height)
    movie_sid = None
    if idx == DEMO_IDX:
        mv = s.shapes.add_movie(
            os.path.join(DOCS, "demo_video.mp4"),
            Inches(vid["x"]), Inches(vid["y"]), Inches(vid["w"]), Inches(vid["h"]),
            poster_frame_image=os.path.join(DOCS, "demo_poster.png"),
            mime_type="video/mp4",
        )
        movie_sid = mv.shape_id
        print("  + nhúng video demo vào slide 6")
    pic_shapes.append((s, movie_sid))

# 4) transition fade + fade-in ảnh
CHILD = """<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="3" fill="hold">
 <p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
  <p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
   <p:par><p:cTn id="5" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect">
    <p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
     <p:set><p:cBhvr><p:cTn id="6" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
      <p:tgtEl><p:spTgt spid="__SID__"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>
      <p:to><p:strVal val="visible"/></p:to></p:set>
     <p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="7" dur="500" fill="hold"/>
      <p:tgtEl><p:spTgt spid="__SID__"/></p:tgtEl><p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst></p:cBhvr>
      <p:tavLst><p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
       <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav></p:tavLst></p:anim>
    </p:childTnLst></p:cTn></p:par>
  </p:childTnLst></p:cTn></p:par>
 </p:childTnLst></p:cTn></p:par>"""
TIMING = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
 <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
  <p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>__C__</p:childTnLst></p:cTn>
  <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
  <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
  </p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""

for s, _sid in pic_shapes:
    el = s._element
    for ex in el.findall("p:transition", NS): el.remove(ex)
    tr = etree.SubElement(el, "{%s}transition" % NS["p"]); tr.set("spd", "med")
    etree.SubElement(tr, "{%s}fade" % NS["p"])
    # fade-in ảnh nền (shape đầu tiên)
    pic_sid = s.shapes[0].shape_id
    for ex in el.findall("p:timing", NS): el.remove(ex)
    el.append(etree.fromstring(TIMING.replace("__C__", CHILD.replace("__SID__", str(pic_sid)))))

prs.save(PPT)
print(f"Saved -> {PPT} ({len(prs.slides._sldIdLst)} slides ảnh + video)")
