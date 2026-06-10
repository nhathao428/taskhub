# -*- coding: utf-8 -*-
"""Dựng THUYET_TRINH_DO_AN.pptx — file thuyết trình NATIVE (chữ sửa được).

Slide PowerPoint gốc: nền tối, shape + textbox editable, tông màu TaskHub.
Kèm transition fade + animation fade-in cho nội dung mỗi slide.
"""
import os, shutil
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

DOCS = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPT = os.path.join(DOCS, "THUYET_TRINH_DO_AN.pptx")

# ---- palette ----
BG       = RGBColor(0x0B, 0x10, 0x20)
BG2      = RGBColor(0x12, 0x1B, 0x36)
CARD     = RGBColor(0x16, 0x20, 0x3E)
CARD2    = RGBColor(0x1B, 0x27, 0x4C)
LINE     = RGBColor(0x2A, 0x36, 0x5E)
INK      = RGBColor(0xEE, 0xF2, 0xFF)
MUTED    = RGBColor(0x9A, 0xA6, 0xC7)
ACCENT   = RGBColor(0x5B, 0x8C, 0xFF)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)
PURPLE   = RGBColor(0xA7, 0x8B, 0xFA)
GREEN    = RGBColor(0x7D, 0xF2, 0xC0)

FONT_H = "Segoe UI Semibold"
FONT_B = "Segoe UI"
FONT_M = "Consolas"

NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

# ---- backup ----
if os.path.exists(PPT):
    stamp = datetime.now().strftime("%Y%m%d-%H%M")
    bak = os.path.join(DOCS, f"THUYET_TRINH_DO_AN.backup-{stamp}.pptx")
    shutil.copy(PPT, bak); print("Backup ->", os.path.basename(bak))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def shape(s, kind, l, t, w, h, fill=None, line=None, line_w=1.0, radius=0.06):
    sp = s.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def text(s, l, t, w, h, runs, size=18, color=INK, bold=False, font=FONT_B,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_sp=1.0):
    """runs: str hoặc list[ (text,{opts}) ] cho nhiều style trong 1 đoạn.
    Cho nhiều đoạn: truyền list các (mục trên)."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.line_spacing = line_sp
        chunks = para if isinstance(para, list) else [(para, {})]
        for txt, o in chunks:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(o.get("size", size)); r.font.bold = o.get("bold", bold)
            r.font.name = o.get("font", font); r.font.color.rgb = o.get("color", color)
            sp = o.get("spacing")
            if sp is not None:
                r.font._rPr.set("spc", str(int(sp * 100)))
    return tb

def eyebrow(s, l, t, label, color=CYAN):
    shape(s, MSO_SHAPE.OVAL, l, t + 0.05, 0.12, 0.12, fill=color)
    text(s, l + 0.22, t - 0.04, 6, 0.4, label.upper(), size=12.5, color=MUTED,
         font=FONT_M, bold=True)

def header(s, num, total=7):
    # wordmark
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, 0.5, 0.42, 0.42, fill=ACCENT, radius=0.28)
    text(s, 0.62, 0.5, 0.42, 0.42, "✓", size=18, color=INK, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.12, 0.5, 3, 0.45, [[("Task", {"bold": True}), ("Hub", {"bold": True, "color": CYAN})]],
         size=19, color=INK, font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)
    text(s, SW.inches - 2.3, 0.52, 1.7, 0.4, f"0{num} / 0{total}", size=12, color=MUTED,
         font=FONT_M, align=PP_ALIGN.RIGHT, bold=True)

def title(s, t_top, parts, size=40):
    text(s, 0.62, t_top, 12, 1.2, [parts], size=size, color=INK, font=FONT_H, bold=True, line_sp=1.0)

def footer(s, txt):
    shape(s, MSO_SHAPE.RECTANGLE, 0.62, SH.inches - 0.62, 12.1, 0.012, fill=LINE)
    text(s, 0.62, SH.inches - 0.55, 12.1, 0.4, txt, size=11.5, color=MUTED, font=FONT_M, bold=True)

def card_feature(s, l, t, w, h, icon, ttl, desc, accent=ACCENT):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill=CARD, line=LINE, radius=0.07)
    text(s, l + 0.32, t + 0.28, w - 0.6, 0.6, icon, size=26)
    text(s, l + 0.32, t + 0.92, w - 0.6, 0.5, ttl, size=18.5, color=INK, font=FONT_H, bold=True)
    text(s, l + 0.32, t + 1.42, w - 0.62, h - 1.5, desc, size=13.5, color=MUTED, line_sp=1.05)

# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = slide()
# decorative glow bands
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 8.7, -1.5, 7, 7, fill=BG2, radius=0.5)
shape(s, MSO_SHAPE.OVAL, 9.6, 0.6, 3.2, 3.2, fill=CARD2)
# eyebrow + brand
eyebrow(s, 0.9, 1.5, "Giới thiệu sản phẩm · Đồ án cơ sở")
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 2.0, 1.05, 1.05, fill=ACCENT, radius=0.26)
text(s, 0.9, 2.0, 1.05, 1.05, "✓", size=40, color=INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 2.15, 2.0, 9, 1.1, [[("Task", {"bold": True}), ("Hub", {"bold": True, "color": CYAN})]],
     size=76, color=INK, font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.9, 3.35, 11.5, 1.2,
     [[("Hệ thống ", {}), ("Quản lý Công việc", {"color": CYAN, "bold": True}),
       (" cho Doanh nghiệp Nhỏ Đa ngành · Tích hợp AI", {})]],
     size=26, color=INK, font=FONT_H, line_sp=1.1)
text(s, 0.9, 4.5, 10.8, 1.0,
     "Nền tảng hợp nhất quản lý nhân sự, chấm công, dự án và tiến độ công việc — kèm trợ lý AI tự động gợi ý nhân viên phù hợp cho từng nhiệm vụ.",
     size=16, color=MUTED, line_sp=1.3)
# author bar
shape(s, MSO_SHAPE.RECTANGLE, 0.9, 5.95, 11.5, 0.012, fill=LINE)
text(s, 0.9, 6.1, 12, 0.9,
     [[("GVHD: ", {"color": MUTED}), ("ThS. Dương Thành Phết", {"bold": True})],
      [("SVTH: ", {"color": MUTED}), ("Nguyễn Nhật Hào", {"bold": True}),
       ("   ·   MSSV 2380612688   ·   Lớp 23DTHC1", {"color": MUTED})],
      [("Trường ĐH Công nghệ TP.HCM · Khoa CNTT · TP.HCM, 5/2026", {"color": MUTED, "size": 12})]],
     size=15, color=INK, line_sp=1.2)

# ============================================================
# SLIDE 2 — BÀI TOÁN & GIẢI PHÁP
# ============================================================
s = slide(); header(s, 2)
eyebrow(s, 0.62, 1.5, "Bài toán")
title(s, 1.85, [("Doanh nghiệp nhỏ đang ", {}), ("vận hành công việc", {"color": CYAN}), (" ra sao?", {})], size=34)
pains = [("📑", "Dữ liệu rời rạc", "Excel, Zalo, giấy tờ — thông tin nằm rải rác, khó tổng hợp."),
         ("🎯", "Phân công cảm tính", "Giao việc theo thói quen, không dựa trên năng lực."),
         ("⏱️", "Chấm công thủ công", "Theo dõi giờ giấc bằng tay, dễ sai sót, khó kiểm soát."),
         ("📉", "Mất dấu tiến độ", "Không rõ dự án tới đâu, ai trễ hạn, hiệu suất ra sao.")]
x0, y0, cw, ch, gx, gy = 0.62, 3.05, 3.55, 1.75, 0.25, 0.22
for i, (ic, t_, d_) in enumerate(pains):
    cx = x0 + (i % 2) * (cw + gx); cy = y0 + (i // 2) * (ch + gy)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch, fill=CARD, line=LINE, radius=0.08)
    text(s, cx + 0.28, cy + 0.25, 0.8, 0.5, ic, size=22)
    text(s, cx + 0.95, cy + 0.28, cw - 1.1, 0.4, t_, size=15.5, color=INK, font=FONT_H, bold=True)
    text(s, cx + 0.95, cy + 0.72, cw - 1.15, ch - 0.8, d_, size=12, color=MUTED, line_sp=1.05)
# solution panel
px, pw = 8.3, 4.45
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, px, 3.05, pw, 3.69, fill=CARD2, line=ACCENT, radius=0.05)
text(s, px + 0.4, 3.35, pw - 0.8, 0.4, "GIẢI PHÁP", size=13, color=CYAN, font=FONT_M, bold=True)
text(s, px + 0.4, 3.75, pw - 0.8, 1.1, [[("Một nền tảng", {})], [("thay cho tất cả.", {})]],
     size=30, color=INK, font=FONT_H, bold=True, line_sp=1.0)
text(s, px + 0.4, 5.05, pw - 0.8, 1.3,
     [[("TaskHub gom nhân sự, chấm công, dự án và công việc về một nơi — thêm lớp ", {}),
       ("AI", {"color": CYAN, "bold": True}), (" giúp phân công đúng người, đúng việc, đúng lúc.", {})]],
     size=14.5, color=INK, line_sp=1.25)
for i, tag in enumerate(["Tập trung", "Tự động hóa", "Realtime"]):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, px + 0.4 + i * 1.35, 6.25, 1.25, 0.4, fill=CARD, line=LINE, radius=0.5)
    text(s, px + 0.4 + i * 1.35, 6.25, 1.25, 0.4, tag, size=11.5, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, "TASKHUB · BÀI TOÁN & GIẢI PHÁP")

# ============================================================
# SLIDE 3 — TÍNH NĂNG
# ============================================================
s = slide(); header(s, 3)
eyebrow(s, 0.62, 1.5, "Tính năng")
title(s, 1.85, [("Mọi thứ để ", {}), ("điều phối công việc", {"color": CYAN})], size=34)
feats = [("👥", "Quản lý nhân sự", "CRUD nhân viên, hồ sơ phòng ban / chức vụ / nhóm."),
         ("📋", "Dự án & công việc", "Tạo dự án, gắn task, theo dõi trạng thái & deadline."),
         ("📍", "Chấm công + GPS", "Check-in/out kèm geofence, duyệt công, lịch sử."),
         ("📊", "Dashboard & báo cáo", "Biểu đồ Chart.js: tiến độ, hiệu suất, đúng hạn."),
         ("🔐", "Phân quyền JWT", "3 vai trò Admin / Manager / Employee."),
         ("📱", "Đa nền tảng", "Web React + mobile Flutter dùng chung REST API.")]
fx, fy, fw, fh, gx, gy = 0.62, 3.0, 3.85, 1.72, 0.25, 0.22
for i, (ic, t_, d_) in enumerate(feats):
    cx = fx + (i % 3) * (fw + gx); cy = fy + (i // 3) * (fh + gy)
    card_feature(s, cx, cy, fw, fh, ic, t_, d_)
footer(s, "TASKHUB · TÍNH NĂNG CHÍNH")

# ============================================================
# SLIDE 4 — AI
# ============================================================
s = slide(); header(s, 4)
eyebrow(s, 0.62, 1.5, "Điểm nhấn · AI")
title(s, 1.85, [("Trợ lý AI ", {}), ("gợi ý phân công", {"color": CYAN})], size=34)
text(s, 0.62, 2.7, 6.5, 0.9, "Không còn giao việc cảm tính — TaskHub dùng Google Gemini để đề xuất nhân viên phù hợp nhất cho mỗi nhiệm vụ.",
     size=15, color=MUTED, line_sp=1.25)
steps = [("1", "Đọc bối cảnh", "Tổng hợp kỹ năng, khối lượng việc & lịch sử hiệu suất của nhân viên."),
         ("2", "Hỏi Gemini", "Gửi yêu cầu + dữ liệu nhân sự lên Gemini qua backend (API key giấu kín)."),
         ("3", "Gợi ý xếp hạng", "Trả về danh sách nhân viên phù hợp kèm điểm số & lý do.")]
sy = 3.75
for n, t_, d_ in steps:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, sy, 0.6, 0.6, fill=ACCENT, radius=0.25)
    text(s, 0.62, sy, 0.6, 0.6, n, size=22, color=INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.4, sy - 0.02, 5.7, 0.4, t_, size=17, color=INK, font=FONT_H, bold=True)
    text(s, 1.4, sy + 0.42, 5.7, 0.7, d_, size=13, color=MUTED, line_sp=1.05)
    sy += 1.0
# mock panel
px, pw = 7.45, 5.25
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, px, 2.55, pw, 4.2, fill=CARD2, line=LINE, radius=0.045)
text(s, px + 0.4, 2.85, 3.2, 0.4, "Gợi ý cho task mới", size=18, color=INK, font=FONT_H, bold=True)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, px + pw - 1.65, 2.85, 1.25, 0.4, fill=CYAN, radius=0.3)
text(s, px + pw - 1.65, 2.85, 1.25, 0.4, "GEMINI 2.5", size=10.5, color=BG, font=FONT_M, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, px + 0.4, 3.45, pw - 0.8, 0.3, "NHIỆM VỤ", size=11, color=MUTED, font=FONT_M, bold=True)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, px + 0.4, 3.78, pw - 0.8, 0.55, fill=CARD, line=LINE, radius=0.18)
text(s, px + 0.6, 3.78, pw - 1.1, 0.55, "“Thiết kế UI màn hình chấm công có bản đồ GPS”", size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
text(s, px + 0.4, 4.5, pw - 0.8, 0.3, "NHÂN VIÊN ĐƯỢC ĐỀ XUẤT", size=11, color=MUTED, font=FONT_M, bold=True)
people = [("Nguyễn Văn A", "Frontend · React", "96%", ACCENT),
          ("Trần Thị B", "Mobile · Flutter", "88%", PURPLE),
          ("Lê Văn C", "Fullstack", "74%", CYAN)]
ry = 4.85
for nm, role, sc, col in people:
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, px + 0.4, ry, pw - 0.8, 0.55, fill=CARD, line=LINE, radius=0.18)
    shape(s, MSO_SHAPE.OVAL, px + 0.55, ry + 0.11, 0.33, 0.33, fill=col)
    text(s, px + 1.0, ry + 0.05, 2.7, 0.3, nm, size=13.5, color=INK, bold=True)
    text(s, px + 1.0, ry + 0.3, 2.7, 0.25, role, size=10.5, color=MUTED)
    text(s, px + pw - 1.2, ry, 0.8, 0.55, sc, size=16, color=CYAN, font=FONT_M, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    ry += 0.62
footer(s, "TASKHUB · TRỢ LÝ AI (GOOGLE GEMINI)")

# ============================================================
# SLIDE 5 — KIẾN TRÚC
# ============================================================
s = slide(); header(s, 5)
eyebrow(s, 0.62, 1.5, "Kiến trúc")
title(s, 1.85, [("Kiến trúc ", {}), ("3 tầng", {"color": CYAN}), (" hiện đại", {})], size=34)
tiers = [("Client", CYAN, ["Web (React + Vite + Tailwind)", "Mobile (Flutter)", "REST qua Axios + JWT"], ["React 18", "Flutter 3", "Chart.js"]),
         ("Application", ACCENT, ["Spring Boot REST API (5000)", "JWT + Spring Security", "Module AI gọi Gemini"], ["Spring Boot 3.5", "Java 17", "JWT"]),
         ("Data", PURPLE, ["PostgreSQL 16 (dữ liệu chính)", "Redis 7 (cache + rate limit)", "Docker Compose"], ["PostgreSQL 16", "Redis 7", "Docker"])]
tx, tw, th, gap = 0.62, 3.75, 2.75, 0.42
card_top, chip_y = 2.95, 5.18
for i, (nm, col, items, chips) in enumerate(tiers):
    cx = tx + i * (tw + gap)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_top, tw, th, fill=CARD, line=LINE, radius=0.05)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_top, tw, 0.09, fill=col, radius=0.5)
    text(s, cx + 0.35, card_top + 0.25, tw - 0.6, 0.5, nm, size=21, color=col, font=FONT_H, bold=True)
    text(s, cx + 0.35, card_top + 0.85, tw - 0.6, 1.4, [[("• " + it, {})] for it in items], size=13, color=INK, line_sp=1.15, sp_after=6)
    cxp = cx + 0.35
    for ch in chips:
        w = 0.18 + len(ch) * 0.085
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cxp, chip_y, w, 0.36, fill=CARD2, line=LINE, radius=0.25)
        text(s, cxp, chip_y, w, 0.36, ch, size=10.5, color=INK, font=FONT_M, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cxp += w + 0.12
    if i < 2:
        text(s, cx + tw - 0.02, card_top + 1.05, gap, 0.6, "→", size=26, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# AI note
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, 6.0, 12.1, 0.72, fill=CARD2, line=LINE, radius=0.18)
text(s, 0.95, 6.0, 0.6, 0.72, "🤖", size=22, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.6, 6.0, 11, 0.72,
     [[("Module AI ở tầng Application — gọi ", {}), ("Google Gemini (gemini-2.5-flash)", {"color": CYAN, "bold": True}),
       (" để gợi ý nhân viên. API key chỉ ở backend, không lộ ra client.", {})]],
     size=13.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.1)
footer(s, "TASKHUB · KIẾN TRÚC & CÔNG NGHỆ")

# ============================================================
# SLIDE 6 — DEMO SẢN PHẨM (video nhúng)
# ============================================================
s = slide(); header(s, 6)
eyebrow(s, 0.62, 1.5, "Demo sản phẩm")
title(s, 1.85, [("Xem ", {}), ("TaskHub", {"color": CYAN}), (" vận hành", {})], size=34)
# left info
text(s, 0.62, 2.95, 3.3, 3.4,
     [[("▶  Đăng nhập & phân quyền", {"bold": True, "color": INK})],
      [("Đăng nhập JWT, vào dashboard theo vai trò.", {"color": MUTED, "size": 12})],
      [(" ", {"size": 6})],
      [("▶  Quản lý dự án & công việc", {"bold": True, "color": INK})],
      [("Tạo dự án, giao task, theo dõi tiến độ.", {"color": MUTED, "size": 12})],
      [(" ", {"size": 6})],
      [("▶  AI gợi ý nhân viên", {"bold": True, "color": INK})],
      [("Chọn task → nhận đề xuất từ Gemini.", {"color": MUTED, "size": 12})],
      [(" ", {"size": 6})],
      [("▶  Chấm công GPS & báo cáo", {"bold": True, "color": INK})],
      [("Check-in geofence, dashboard biểu đồ.", {"color": MUTED, "size": 12})]],
     size=15, color=INK, line_sp=1.15, sp_after=2)
# video frame (browser-style card)
vx, vy, vw, vh = 4.25, 2.75, 8.45, 4.05
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, vx, vy, vw, vh, fill=CARD2, line=LINE, radius=0.04)
# title bar dots
for i, col in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
    shape(s, MSO_SHAPE.OVAL, vx + 0.3 + i * 0.28, vy + 0.25, 0.16, 0.16, fill=col)
text(s, vx + 1.4, vy + 0.16, 4, 0.35, "taskhub · demo", size=12, color=MUTED, font=FONT_M, anchor=MSO_ANCHOR.MIDDLE)
# embedded movie with poster
mv_x, mv_y, mv_w, mv_h = vx + 0.18, vy + 0.62, vw - 0.36, vh - 0.8
try:
    s.shapes.add_movie(os.path.join(DOCS, "demo_video.mp4"),
                       Inches(mv_x), Inches(mv_y), Inches(mv_w), Inches(mv_h),
                       poster_frame_image=os.path.join(DOCS, "demo_poster.png"),
                       mime_type="video/mp4")
    print("  + embedded demo_video.mp4")
except Exception as e:
    print("  ! movie embed failed:", e)
footer(s, "TASKHUB · DEMO SẢN PHẨM  ·  bấm ▶ để phát video")

# ============================================================
# SLIDE 7 — ĐỐI TƯỢNG
# ============================================================
s = slide(); header(s, 7)
eyebrow(s, 0.62, 1.5, "Đối tượng sử dụng")
title(s, 1.85, [("Dành cho ", {}), ("ai?", {"color": CYAN})], size=34)
who = [("🏪", "Doanh nghiệp nhỏ đa ngành", "Cần công cụ quản trị công việc gọn nhẹ, không cồng kềnh như ERP."),
       ("👔", "Quản lý / chủ doanh nghiệp", "Muốn nhìn toàn cảnh nhân sự & tiến độ, phân công nhanh, có cơ sở."),
       ("🧑‍💻", "Nhân viên", "Tự xem việc được giao, cập nhật trạng thái, tự chấm công minh bạch.")]
wx, wy, ww, wh, g = 0.62, 3.0, 3.93, 2.0, 0.25
for i, (ic, t_, d_) in enumerate(who):
    cx = wx + i * (ww + g)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, wy, ww, wh, fill=CARD, line=LINE, radius=0.07)
    text(s, cx + 0.35, wy + 0.3, 1, 0.6, ic, size=30)
    text(s, cx + 0.35, wy + 1.0, ww - 0.7, 0.5, t_, size=16.5, color=INK, font=FONT_H, bold=True)
    text(s, cx + 0.35, wy + 1.45, ww - 0.7, wh - 1.5, d_, size=13, color=MUTED, line_sp=1.1)
# closing band
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, 5.45, 12.1, 1.45, fill=CARD2, line=ACCENT, radius=0.07)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.95, 5.78, 0.8, 0.8, fill=ACCENT, radius=0.26)
text(s, 0.95, 5.78, 0.8, 0.8, "✓", size=30, color=INK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 2.0, 5.7, 7.5, 1.0,
     [[("TaskHub", {"color": CYAN, "bold": True}), (" — Quản lý công việc thông minh hơn.", {"bold": True})],
      [("Đồ án cơ sở · Khoa CNTT · ĐH Công nghệ TP.HCM", {"color": MUTED, "size": 12})]],
     size=22, color=INK, font=FONT_H, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.15)
text(s, 9.4, 5.65, 3.1, 1.15,
     [[("GVHD: ThS. Dương Thành Phết", {})], [("SVTH: Nguyễn Nhật Hào", {})],
      [("MSSV 2380612688 · 23DTHC1", {"color": MUTED, "size": 11})]],
     size=12.5, color=INK, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.2)
footer(s, "TASKHUB · CẢM ƠN ĐÃ LẮNG NGHE")

# ============================================================
# Transition fade + animation fade-in nội dung
# ============================================================
TIMING = """
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
 <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
  <p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
   __CHILDREN__
  </p:childTnLst></p:cTn>
  <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
  <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
  </p:seq>
 </p:childTnLst></p:cTn></p:par></p:tnLst>
</p:timing>"""

CHILD = """
<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="__A__" fill="hold">
 <p:stCondLst><p:cond delay="__DELAY__"/></p:stCondLst><p:childTnLst>
  <p:par><p:cTn id="__B__" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
   <p:par><p:cTn id="__C__" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect">
    <p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
     <p:set><p:cBhvr><p:cTn id="__D__" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
      <p:tgtEl><p:spTgt spid="__SID__"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>
      <p:to><p:strVal val="visible"/></p:to></p:set>
     <p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="__E__" dur="450" fill="hold"/>
      <p:tgtEl><p:spTgt spid="__SID__"/></p:tgtEl><p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst></p:cBhvr>
      <p:tavLst><p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
       <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav></p:tavLst></p:anim>
    </p:childTnLst></p:cTn></p:par>
  </p:childTnLst></p:cTn></p:par>
 </p:childTnLst></p:cTn></p:par>"""

for s in prs.slides:
    el = s._element
    # transition
    for ex in el.findall("p:transition", NS): el.remove(ex)
    tr = etree.SubElement(el, "{%s}transition" % NS["p"]); tr.set("spd", "med")
    etree.SubElement(tr, "{%s}fade" % NS["p"])
    # animations: stagger các shape có text
    kids = []; base = 10; delay = 0
    for sh in s.shapes:
        if not sh.has_text_frame: continue
        if not sh.text_frame.text.strip(): continue
        ids = [base, base + 1, base + 2, base + 3, base + 4]; base += 5
        b = CHILD.replace("__SID__", str(sh.shape_id)).replace("__DELAY__", str(delay))
        for tag, val in zip(["__A__", "__B__", "__C__", "__D__", "__E__"], ids):
            b = b.replace(tag, str(val))
        kids.append(b); delay += 80
    if kids:
        for ex in el.findall("p:timing", NS): el.remove(ex)
        el.append(etree.fromstring(TIMING.replace("__CHILDREN__", "\n".join(kids))))

prs.save(PPT)
print(f"Saved native deck -> {PPT} ({len(prs.slides._sldIdLst)} slides)")
